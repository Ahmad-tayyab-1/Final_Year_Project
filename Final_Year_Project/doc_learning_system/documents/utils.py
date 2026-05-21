import fitz
import re
import os
import json
import requests
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.units import inch
from groq import Groq
from .youtube import extract_youtube_video_id


class DocumentProcessor:

    def extract_text_with_pages(self, file_path):
        doc = fitz.open(file_path)
        pages = []
        full_text = ""
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            pages.append({"page": page_num, "text": text})
            full_text += f"\\n--- Page {page_num} ---\\n{text}"
        doc.close()
        return full_text, pages, len(pages)

    def extract_text_from_docx(self, file_path):
        doc = DocxDocument(file_path)
        return "\\n".join([p.text for p in doc.paragraphs])

    def extract_text_from_doc(self, file_path):
        """Extract text from older .doc files (Office 97-2003 format)."""
        try:
            # Try using python-docx first (works with some .doc files)
            doc = DocxDocument(file_path)
            return "\\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            # If that fails, try to convert .doc to .docx using LibreOffice or similar
            try:
                import subprocess
                import tempfile

                # Create a temporary directory for conversion
                with tempfile.TemporaryDirectory() as tmp_dir:
                    output_path = os.path.join(tmp_dir, "converted.docx")

                    # Try using LibreOffice command line conversion
                    try:
                        subprocess.run([
                            'soffice', '--headless', '--convert-to', 'docx',
                            '--outdir', tmp_dir, file_path
                        ], check=True, capture_output=True, timeout=30)

                        converted_file = os.path.join(tmp_dir, os.path.splitext(
                            os.path.basename(file_path))[0] + ".docx")
                        if os.path.exists(converted_file):
                            doc = DocxDocument(converted_file)
                            return "\\n".join([p.text for p in doc.paragraphs])
                    except (subprocess.CalledProcessError, FileNotFoundError):
                        pass
            except Exception:
                pass

            # Fallback: Extract text from binary format
            try:
                from docx.oxml import parse_xml
                from zipfile import ZipFile

                # Try to read as if it's a zipped XML (some .doc files are actually docx)
                try:
                    with ZipFile(file_path, 'r') as zip_ref:
                        xml_content = zip_ref.read('word/document.xml')
                        import re
                        # Remove XML tags
                        text = re.sub(r'<[^>]+>', '',
                                      xml_content.decode('utf-8'))
                        return text.strip()
                except:
                    pass
            except Exception:
                pass

            # Final fallback: Extract ASCII text from binary
            try:
                with open(file_path, 'rb') as f:
                    content = f.read()
                    # Try to decode common text patterns in .doc files
                    text = []
                    current = b''
                    for byte in content:
                        if 32 <= byte <= 126:  # Printable ASCII
                            current += bytes([byte])
                        else:
                            if len(current) > 4:
                                try:
                                    text.append(current.decode(
                                        'utf-8', errors='ignore'))
                                except:
                                    pass
                            current = b''

                    result = ' '.join(text)
                    return result if result.strip() else "Could not extract text from .doc file. File may be corrupted."
            except Exception as extract_err:
                return f"Error processing .doc file: {str(extract_err)}"

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\\n"
        return text

    def extract_text_from_xlsx(self, file_path):
        """Extract text from XLSX file."""
        try:
            import openpyxl
        except ImportError:
            return "openpyxl not installed for XLSX processing"

        wb = openpyxl.load_workbook(file_path)
        text = ""

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"Sheet: {sheet_name}\\n"

            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell)
                                    for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\\n"

        return text

    def extract_text_from_txt(self, file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def clean_text(self, text):
        return re.sub(r"\\s+", " ", text).strip()

    def generate_office_html(self, doc):
        """Generate HTML representation of Office documents."""
        file_path = doc.file.path
        file_type = doc.file_type

        if file_type in ('docx', 'doc'):
            return self._docx_to_html(file_path)
        elif file_type == 'pptx':
            return self._pptx_to_html(file_path)
        elif file_type == 'xlsx':
            return self._xlsx_to_html(file_path)
        else:
            return f"<p>Unsupported file type: {file_type}</p>"

    def convert_office_to_pdf(self, doc):
        """Convert Office document to PDF using LibreOffice and return PDF file path."""
        file_path = doc.file.path
        file_type = doc.file_type

        # All office types are handled by LibreOffice
        if file_type in ['doc', 'docx', 'pptx', 'xlsx']:
            pdf_path = file_path.rsplit('.', 1)[0] + '.pdf'
            self._libreoffice_to_pdf(file_path, pdf_path)
        else:
            raise ValueError(f"Unsupported conversion for {file_type}")

        return pdf_path

    def _libreoffice_to_pdf(self, file_path, pdf_path):
        """Convert any Office file to PDF using LibreOffice headless mode."""
        import subprocess
        import shutil
        import tempfile

        # LibreOffice executable – try PATH first, then the default install location
        soffice_candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "soffice",  # if it is on PATH
        ]
        soffice = None
        for candidate in soffice_candidates:
            if os.path.exists(candidate) if os.sep in candidate else True:
                soffice = candidate
                break
        if soffice is None:
            raise RuntimeError(
                "LibreOffice (soffice) not found. Please install LibreOffice.")

        # LibreOffice writes the output PDF next to the source file inside --outdir.
        # We work in a temp directory to avoid polluting the media folder during conversion.
        with tempfile.TemporaryDirectory() as tmp_dir:
            # Copy the source document into the temp dir
            tmp_src = os.path.join(tmp_dir, os.path.basename(file_path))
            shutil.copy2(file_path, tmp_src)

            # Run the conversion
            result = subprocess.run(
                [soffice, '--headless', '--convert-to',
                    'pdf', '--outdir', tmp_dir, tmp_src],
                capture_output=True,
                timeout=120,
            )

            if result.returncode != 0:
                stderr = result.stderr.decode(errors='ignore')
                raise RuntimeError(
                    f"LibreOffice PDF conversion failed: {stderr}")

            # LibreOffice names the output <original_stem>.pdf
            base_stem = os.path.splitext(os.path.basename(file_path))[0]
            tmp_pdf = os.path.join(tmp_dir, base_stem + '.pdf')

            if not os.path.exists(tmp_pdf):
                raise RuntimeError(
                    "LibreOffice ran but did not produce a PDF. "
                    "Check that the file is not password-protected."
                )

            # Move the finished PDF to the target location
            shutil.move(tmp_pdf, pdf_path)

    def _docx_to_html(self, file_path):
        """Convert DOCX to full-fidelity HTML using mammoth (preserves images, tables, formatting)."""
        import mammoth
        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(f)
        return result.value

    def _pptx_to_html(self, file_path):
        """Convert PPTX to HTML slides with ALL images extracted."""
        import base64
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        prs = Presentation(file_path)
        html_parts = ["<div class='pptx-viewer'>"]

        def _extract_shapes(shapes, parts_list):
            """Recursively extract content from shapes including groups."""
            for shape in shapes:
                # Recurse into group shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        _extract_shapes(shape.shapes, parts_list)
                    except Exception:
                        pass
                    continue

                # Try to get image from shape (pictures, placeholders with images)
                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = image.content_type
                    b64 = base64.b64encode(img_bytes).decode('utf-8')
                    parts_list.append(
                        f"<div style='text-align:center; margin:12px 0;'>"
                        f"<img src='data:{content_type};base64,{b64}' "
                        f"style='max-width:100%; height:auto; border-radius:4px;' />"
                        f"</div>"
                    )
                    # Also show text if shape has both image and text
                    if hasattr(shape, "text") and shape.text.strip():
                        parts_list.append(f"<p>{shape.text}</p>")
                    continue
                except Exception:
                    pass

                # Try to get image from shape fill (background fill images)
                try:
                    fill = shape.fill
                    if fill and fill.type is not None:
                        blip = fill._fill.findall(
                            './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                        for b in blip:
                            rId = b.get(
                                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if rId:
                                rel = shape.part.rels[rId]
                                img_bytes = rel.target_part.blob
                                content_type = rel.target_part.content_type
                                b64_data = base64.b64encode(
                                    img_bytes).decode('utf-8')
                                parts_list.append(
                                    f"<div style='text-align:center; margin:12px 0;'>"
                                    f"<img src='data:{content_type};base64,{b64_data}' "
                                    f"style='max-width:100%; height:auto; border-radius:4px;' />"
                                    f"</div>"
                                )
                except Exception:
                    pass

                # Handle tables
                if shape.has_table:
                    table = shape.table
                    parts_list.append("<table>")
                    for row_idx, row in enumerate(table.rows):
                        parts_list.append("<tr>")
                        tag = "th" if row_idx == 0 else "td"
                        for cell in row.cells:
                            parts_list.append(f"<{tag}>{cell.text}</{tag}>")
                        parts_list.append("</tr>")
                    parts_list.append("</table>")
                # Handle text
                elif hasattr(shape, "text") and shape.text.strip():
                    parts_list.append(f"<p>{shape.text}</p>")

        for i, slide in enumerate(prs.slides, 1):
            html_parts.append(f"<div class='slide' id='slide-{i}'>")
            html_parts.append(f"<h2>Slide {i}</h2>")

            # Extract slide background image
            try:
                bg = slide.background
                bg_fill = bg.fill
                if bg_fill._fill is not None:
                    blips = bg_fill._fill.findall(
                        './/{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    for b in blips:
                        rId = b.get(
                            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                        if rId:
                            rel = slide.part.rels[rId]
                            img_bytes = rel.target_part.blob
                            content_type = rel.target_part.content_type
                            b64_data = base64.b64encode(
                                img_bytes).decode('utf-8')
                            html_parts.append(
                                f"<div style='text-align:center; margin:12px 0;'>"
                                f"<img src='data:{content_type};base64,{b64_data}' "
                                f"style='max-width:100%; height:auto; border-radius:4px; "
                                f"box-shadow: 0 2px 8px rgba(0,0,0,0.1);' />"
                                f"</div>"
                            )
            except Exception:
                pass

            # Extract all shapes recursively
            _extract_shapes(slide.shapes, html_parts)

            html_parts.append("</div>")

        html_parts.append("</div>")
        return "\n".join(html_parts)

    def _xlsx_to_html(self, file_path):
        """Convert XLSX to HTML table."""
        try:
            import openpyxl
        except ImportError:
            return "<p>openpyxl not installed. Run: pip install openpyxl</p>"

        wb = openpyxl.load_workbook(file_path)
        html_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            html_parts.append(f"<h2>{sheet_name}</h2>")
            html_parts.append(
                "<table border='1' style='border-collapse: collapse;'>")

            for row in sheet.iter_rows(values_only=True):
                html_parts.append("<tr>")
                for cell in row:
                    if cell is not None:
                        html_parts.append(
                            f"<td style='padding: 5px;'>{cell}</td>")
                    else:
                        html_parts.append("<td style='padding: 5px;'></td>")
                html_parts.append("</tr>")

            html_parts.append("</table>")

        return "\n".join(html_parts)

    # NOTE: _docx_to_pdf / _pptx_to_pdf / _xlsx_to_pdf are no longer used.
    # All office-to-PDF conversion now goes through _libreoffice_to_pdf above.
    # Kept as a commented-out fallback reference only.

    def get_relevant_context(self, query, document_text, max_chars=4000):
        """Return the most query-relevant portion of the document."""
        query_words = set(re.findall(r"\\w+", query.lower()))
        if not query_words:
            return document_text[:max_chars]

        chunks = re.split(r"--- Page \\d+ ---", document_text)
        chunks = [c.strip() for c in chunks if c.strip()]
        if not chunks:
            return document_text[:max_chars]

        scored = sorted(
            [(len(set(re.findall(r"\\w+", c.lower())) & query_words), c)
             for c in chunks],
            key=lambda x: x[0], reverse=True
        )

        context = ""
        for _, chunk in scored:
            if len(context) + len(chunk) > max_chars:
                break
            context += chunk + "\\n\\n"
        return context.strip() or document_text[:max_chars]


class YouTubeProcessor:
    """Fetch and process YouTube video transcripts."""

    def extract_video_id(self, url: str):
        return extract_youtube_video_id(url)

    def get_video_metadata(self, url_or_video_id: str):
        """Return public YouTube metadata available without an API key."""
        video_id = self.extract_video_id(url_or_video_id) or url_or_video_id
        if not video_id:
            return {}

        video_url = f"https://www.youtube.com/watch?v={video_id}"
        try:
            resp = requests.get(
                "https://www.youtube.com/oembed",
                params={"url": video_url, "format": "json"},
                timeout=8,
            )
            resp.raise_for_status()
            data = resp.json()
            return {
                "title": data.get("title", ""),
                "author_name": data.get("author_name", ""),
                "author_url": data.get("author_url", ""),
                "video_url": video_url,
            }
        except Exception:
            return {"video_url": video_url}

    def get_transcript(self, video_id: str):
        """
        Returns (transcript_text, video_title_or_None).
        On Render : Invidious instances → Supadata (if key set) → clear error.
        Locally   : youtube-transcript-api directly.
        """
        import logging
        is_render = 'RENDER' in os.environ

        if is_render:
            # ── PRIMARY: Invidious public instances ────────────────────────
            # Full updated list from https://api.invidious.io/instances.json
            INVIDIOUS_INSTANCES = [
                "https://invidious.io",
                "https://inv.nadeko.net",
                "https://invidious.nerdvpn.de",
                "https://iv.ggtyler.dev",
                "https://invidious.incogniweb.net",
                "https://invidious.slipfox.xyz",
                "https://invidious.reallyaweso.me",
                "https://invidious.privacyredirect.com",
                "https://yt.artemislena.eu",
                "https://invidious.darkness.services",
            ]

            last_error = ""
            for base in INVIDIOUS_INSTANCES:
                try:
                    # Step 1: fetch caption track list
                    r = requests.get(
                        f"{base}/api/v1/captions/{video_id}",
                        timeout=8,
                    )
                    if r.status_code != 200:
                        logging.warning(
                            f"[Invidious] {base} → HTTP {r.status_code}")
                        last_error = f"HTTP {r.status_code} from {base}"
                        continue

                    captions = r.json().get("captions", [])
                    if not captions:
                        logging.warning(
                            f"[Invidious] {base} → no captions available")
                        last_error = "No captions found on this video"
                        continue

                    # Step 2: prefer English, fall back to first track
                    track = next(
                        (c for c in captions if c.get(
                            "languageCode", "").startswith("en")),
                        captions[0],
                    )

                    # Step 3: fetch VTT file
                    vtt_url = base + track["url"]
                    vtt_r = requests.get(vtt_url, timeout=10)
                    if vtt_r.status_code != 200:
                        logging.warning(
                            f"[Invidious] {base} VTT fetch → HTTP {vtt_r.status_code}")
                        last_error = f"VTT fetch failed HTTP {vtt_r.status_code}"
                        continue

                    # Step 4: parse VTT
                    transcript_text = self._parse_vtt(vtt_r.text)
                    if not transcript_text.strip():
                        logging.warning(
                            f"[Invidious] {base} → VTT parsed empty")
                        last_error = "VTT file was empty after parsing"
                        continue

                    logging.info(f"[Invidious] SUCCESS via {base}")
                    metadata = self.get_video_metadata(video_id)
                    return transcript_text, metadata.get("title")

                except requests.exceptions.Timeout:
                    logging.warning(f"[Invidious] {base} → timeout")
                    last_error = f"{base} timed out"
                except Exception as e:
                    logging.warning(f"[Invidious] {base} → {e}")
                    last_error = str(e)

            # ── SECONDARY: Supadata (free tier, 100 req/month) ─────────────
            supadata_key = os.environ.get("SUPADATA_API_KEY", "")
            if supadata_key:
                try:
                    logging.info("[Supadata] Trying Supadata API...")
                    r = requests.get(
                        "https://api.supadata.ai/v1/youtube/transcript",
                        params={"videoId": video_id, "lang": "en"},
                        headers={"x-api-key": supadata_key},
                        timeout=15,
                    )
                    r.raise_for_status()
                    content = r.json().get("content", [])
                    if content:
                        lines = []
                        for seg in content:
                            start = seg.get("offset", 0) / 1000  # ms → seconds
                            text = seg.get("text", "").strip()
                            mins, secs = int(start // 60), int(start % 60)
                            if text:
                                lines.append(f"[{mins:02d}:{secs:02d}] {text}")
                        transcript_text = "\n".join(lines)
                        if transcript_text:
                            logging.info("[Supadata] SUCCESS")
                            metadata = self.get_video_metadata(video_id)
                            return transcript_text, metadata.get("title")
                except Exception as e:
                    logging.warning(f"[Supadata] failed: {e}")

            # ── ALL FAILED on Render ────────────────────────────────────────
            # Do NOT fall through to youtube-transcript-api here —
            # it will always be IP-blocked on cloud hosts.
            raise RuntimeError(
                "Could not fetch the transcript on the server. "
                "YouTube blocks cloud IPs and all bypass services are currently unavailable. "
                "Please try again in a few minutes, or try a different video. "
                f"(Last error: {last_error})"
            )

        # ── LOCAL: youtube-transcript-api directly ─────────────────────────
        try:
            from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound
        except ImportError:
            raise RuntimeError("youtube-transcript-api is not installed.")

        api = YouTubeTranscriptApi()
        transcript_list = None

        try:
            transcript_list = api.fetch(video_id) if hasattr(
                api, 'fetch') else YouTubeTranscriptApi.get_transcript(video_id)
        except TranscriptsDisabled:
            raise RuntimeError("This video has captions disabled.")
        except NoTranscriptFound:
            try:
                tl = api.list(video_id) if hasattr(
                    api, 'list') else YouTubeTranscriptApi.list_transcripts(video_id)
                try:
                    transcript_list = tl.find_generated_transcript(
                        ["en"]).fetch()
                except Exception:
                    transcript_list = next(iter(tl)).fetch()
            except Exception:
                raise RuntimeError(
                    "No transcript available for this video. "
                    "Only videos with captions (auto or manual) are supported."
                )
        except Exception as e:
            raise RuntimeError(f"Could not fetch transcript: {e}")

        lines = []
        for entry in transcript_list:
            if isinstance(entry, dict):
                start, text = entry["start"], entry["text"]
            else:
                start, text = entry.start, entry.text
            mins, secs = int(start // 60), int(start % 60)
            lines.append(f"[{mins:02d}:{secs:02d}] {text}")

        metadata = self.get_video_metadata(video_id)
        return "\n".join(lines), metadata.get("title")

    def _parse_vtt(self, vtt_text: str) -> str:
        """Parse a WebVTT string into [MM:SS] timestamped lines."""
        import re
        lines_out = []
        lines = vtt_text.splitlines()
        i = 0
        seen_texts = set()  # VTT files often have duplicate lines near each other

        while i < len(lines):
            line = lines[i].strip()
            # Match a timestamp cue line: 00:00:01.000 --> 00:00:04.000
            ts_match = re.match(
                r"(\d{1,2}:\d{2}:\d{2}\.\d+|\d{2}:\d{2}\.\d+)\s+-->\s+", line
            )
            if ts_match:
                ts_raw = ts_match.group(1)
                # Normalise HH:MM:SS.mmm or MM:SS.mmm → seconds
                parts = ts_raw.split(":")
                if len(parts) == 3:
                    h, m, s = int(parts[0]), int(parts[1]), float(parts[2])
                else:
                    h, m, s = 0, int(parts[0]), float(parts[1])
                total = h * 3600 + m * 60 + s
                mins, secs = int(total // 60), int(total % 60)

                # Collect text lines that follow until blank line
                i += 1
                text_parts = []
                while i < len(lines) and lines[i].strip():
                    # strip VTT inline tags
                    clean = re.sub(r"<[^>]+>", "", lines[i].strip())
                    if clean:
                        text_parts.append(clean)
                    i += 1

                text = " ".join(text_parts).strip()
                if text and text not in seen_texts:
                    seen_texts.add(text)
                    lines_out.append(f"[{mins:02d}:{secs:02d}] {text}")
            else:
                i += 1

        return "\n".join(lines_out)
        '''
    def get_transcript(self, video_id: str):
        """
        Returns (transcript_text, video_title_or_None).
        Raises RuntimeError on failure.
        """
        try:
            from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled, NoTranscriptFound
        except ImportError:
            raise RuntimeError("youtube-transcript-api is not installed.")

        # Build API instance
        if 'RENDER' in os.environ:
            cookies_b64 = os.environ.get('YOUTUBE_COOKIES_B64', '')
            import logging
            logging.warning(f"COOKIES_B64 length: {len(cookies_b64)}")
            if cookies_b64:
               # if 'RENDER' in os.environ:
                #    cookies_b64 = os.environ.get('YOUTUBE_COOKIES_B64', '')
               # if cookies_b64:
                import base64
                import tempfile
                import http.cookiejar
                from requests import Session

                cookies_data = base64.b64decode(
                    cookies_b64).decode('utf-8', errors='ignore')
                tmp = tempfile.NamedTemporaryFile(
                    mode='w', encoding='utf-8', suffix='.txt', delete=False)
                tmp.write(cookies_data)
                tmp.flush()
                tmp.close()

                session = Session()
                cj = http.cookiejar.MozillaCookieJar(tmp.name)
                cj.load(ignore_discard=True, ignore_expires=True)
                cookies_to_remove = []
                for cookie in cj:
                    try:
                        # Test if the cookie's value and name can be encoded to latin-1
                        if cookie.value:
                            cookie.value.encode('latin-1')
                        if cookie.name:
                            cookie.name.encode('latin-1')
                    except UnicodeEncodeError:
                        # If it fails, mark it for removal
                        cookies_to_remove.append(cookie)

                # Safely clear out the problematic cookies from the jar
                for cookie in cookies_to_remove:
                    cj.clear(cookie.domain, cookie.path, cookie.name)
                # === NEWLY ADDED CODE END ===

                session.cookies = cj
                api = YouTubeTranscriptApi(http_client=session)
            else:
                api = YouTubeTranscriptApi()
        else:
            api = YouTubeTranscriptApi()

        transcript_list = None

        try:
            if hasattr(api, 'fetch'):
                transcript_list = api.fetch(video_id)
            else:
                transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
        except TranscriptsDisabled:
            raise RuntimeError("This video has captions disabled.")
        except NoTranscriptFound:
            try:
                if hasattr(api, 'list'):
                    tl = api.list(video_id)
                else:
                    tl = YouTubeTranscriptApi.list_transcripts(video_id)
                try:
                    transcript_list = tl.find_generated_transcript(
                        ["en"]).fetch()
                except Exception:
                    transcript_list = next(iter(tl)).fetch()
            except Exception:
                raise RuntimeError(
                    "No transcript available for this video. "
                    "Only videos with captions (auto or manual) are supported."
                )
        except Exception as e:
            raise RuntimeError(f"Could not fetch transcript: {e}")

        # Build timestamped text
        lines = []
        for entry in transcript_list:
            if isinstance(entry, dict):
                start = entry["start"]
                text = entry["text"]
            else:
                start = entry.start
                text = entry.text
            mins = int(start // 60)
            secs = int(start % 60)
            lines.append(f"[{mins:02d}:{secs:02d}] {text}")

        transcript_text = "\n".join(lines)
        metadata = self.get_video_metadata(video_id)
        return transcript_text, metadata.get("title")
'''


class WebSearch:
    """Small live-web search helper that returns snippets for the AI prompt."""

    def search(self, query, max_results=5):
        try:
            response = requests.get(
                "https://duckduckgo.com/html/",
                params={"q": query},
                headers={"User-Agent": "Mozilla/5.0"},
                timeout=10,
            )
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            results = []

            for result in soup.select(".result"):
                title_el = result.select_one(".result__a")
                snippet_el = result.select_one(".result__snippet")
                if not title_el:
                    continue
                title = title_el.get_text(" ", strip=True)
                url = title_el.get("href", "")
                snippet = snippet_el.get_text(
                    " ", strip=True) if snippet_el else ""
                if title and url:
                    results.append(
                        {"title": title, "url": url, "snippet": snippet})
                if len(results) >= max_results:
                    break

            return results
        except Exception:
            return []

    def format_results(self, results):
        if not results:
            return ""
        lines = []
        for index, item in enumerate(results, 1):
            lines.append(
                f"[{index}] {item['title']}\nURL: {item['url']}\nSnippet: {item.get('snippet', '')}"
            )
        return "\n\n".join(lines)


class AIAssistant:

    def __init__(self):
        self._api_key = os.getenv("GROQ_API_KEY")
        self._client = None

    @property
    def client(self):
        if self._client is None and self._api_key:
            self._client = Groq(api_key=self._api_key)
        return self._client

    def _casual_response(self, query):
        """Return a friendly reply for greetings or tiny non-document messages."""
        normalized = re.sub(r"[^a-z0-9\s]", " ", query.lower()).strip()
        normalized = re.sub(r"\s+", " ", normalized)
        if not normalized:
            return None

        greetings = {
            "hi", "hii", "hello", "hey", "yo", "salam", "assalamualaikum",
            "assalamu alaikum", "good morning", "good afternoon", "good evening",
        }
        check_ins = {
            "how are you", "how r u", "how are u", "what s up", "whats up",
            "wassup", "sup", "how you doing",
        }
        thanks = {"thanks", "thank you", "thx", "ty"}
        vague = {"what", "hwat", "ok", "okay", "hmm", "huh"}

        if normalized in greetings:
            return (
                "Hey! I'm here with you. How are you doing? "
                "You can ask me anything about this document or video whenever you're ready."
            )
        if normalized in check_ins:
            return (
                "I'm doing well, thanks for asking. How are you? "
                "If you want, we can go through this material together."
            )
        if normalized in thanks:
            return "You're welcome. Ask me the next thing you want to understand from this material."
        if normalized in vague:
            return (
                "I'm here. Ask me a specific question about the document or video, "
                "or tell me what topic you want explained."
            )
        return None

    # ── Answer with smart context ────────────────────────────
    def answer_question(self, query, document_text, source_info="", web_results=""):
        casual = self._casual_response(query)
        if casual:
            return casual, None

        if not self.client:
            return "Groq API key not configured.", None
        try:
            proc = DocumentProcessor()
            context = proc.get_relevant_context(query, document_text, 4000)
            resp = self.client.chat.completions.create(
                messages=[
                    {"role": "system",
                     "content": (
                         "You are DocuLearn AI, a natural chatbot and educational assistant. "
                         "First understand what the user is asking, then choose the right mode:\n"
                         "1. For greetings, small talk, thanks, or unclear tiny messages, reply like a normal friendly chatbot.\n"
                         "2. For questions about the uploaded document/video, use the source info and transcript/document context.\n"
                         "3. For questions whose answer is not in the uploaded material, use the live web results if provided. If no web results are provided, answer from your general knowledge when possible and say briefly that it is general knowledge.\n"
                         "4. If you use live web results, say that you checked live web results and include the most relevant source URL(s).\n"
                         "5. Answer the exact question first. Keep answers conversational and concise unless the user asks for a summary or detailed explanation.\n"
                         "6. Do not dump a script-style outline or summarize the whole video/document unless asked.\n"
                         "7. Cite page numbers or timestamps only when useful and identifiable from markers like --- Page N --- or [MM:SS]."
                     )},
                    {"role": "user",
                     "content": (
                         f"Source info:\n{source_info or 'No extra source metadata available.'}\n\n"
                         f"Relevant document/video context:\n{context}\n\n"
                         f"Live web results:\n{web_results or 'No live web results were provided for this question.'}\n\n"
                         f"User question:\n{query}\n\n"
                         "Reply naturally and directly. If source metadata answers the question, use it. "
                         "If live web results are provided and useful, use them. If the context does not contain the answer but you know it generally, answer generally and make that clear."
                     )},
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.6,
                max_tokens=900,
            )
            answer = resp.choices[0].message.content
            page_match = re.search(r"[Pp]age\\s+(\\d+)", answer)
            page_ref = int(page_match.group(1)) if page_match else None
            return answer, page_ref
        except Exception as e:
            return f"Error: {e}", None

    # ── Explain selected text ────────────────────────────────
    def explain_text(self, selected_text, document_context=""):
        if not self.client:
            return "Groq API key not configured."
        try:
            resp = self.client.chat.completions.create(
                messages=[
                    {"role": "system",
                     "content": "Explain the given text clearly and comprehensively. Use examples where helpful and format your reply beautifully using markdown features (like bolding and bullet points)."},
                    {"role": "user",
                     "content": f"Text:\n{selected_text[:2000]}]\n\nContext:\n{document_context[:1200]}\n\nDetailed Explanation:"},
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.7, max_tokens=1200,
            )
            return resp.choices[0].message.content
        except Exception as e:
            return f"Error: {e}"

    # ── Flashcards ────────────────────────────────────────────
    def generate_flashcards(self, text, count=5):
        if not self.client:
            return []
        try:
            resp = self.client.chat.completions.create(
                messages=[
                    {"role": "system",
                     "content": "Create educational flashcards. Respond ONLY with a valid JSON array of objects, each with \"front\" (question) and \"back\" (answer) keys. No markdown, no explanation."},
                    {"role": "user",
                     "content": f"Create exactly {count} flashcards from:\n\n{text[:2500]}"},
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.8, max_tokens=900,
            )
            raw = re.sub(
                r"```[a-z]*", "", resp.choices[0].message.content).strip().strip("`")
            data = json.loads(raw)
            if isinstance(data, list):
                return [{"front": str(d.get("front", "")), "back": str(d.get("back", ""))}
                        for d in data if d.get("front") and d.get("back")][:count]
        except Exception as e:
            print(f"Flashcard error: {e}")
        return []

    # ── MCQs ──────────────────────────────────────────────────
    def generate_mcqs(self, text, count=3):
        if not self.client:
            return []
        try:
            resp = self.client.chat.completions.create(
                messages=[
                    {"role": "system",
                     "content": ('Create MCQs. Respond ONLY with a valid JSON array. Each object must have: '
                                 '"question", "option_a", "option_b", "option_c", "option_d", '
                                 '"correct_answer" (A/B/C/D), "explanation". No markdown.')},
                    {"role": "user",
                     "content": f"Create exactly {count} MCQs from:\\n\\n{text[:2500]}"},
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.7, max_tokens=1200,
            )
            raw = re.sub(
                r"```[a-z]*", "", resp.choices[0].message.content).strip().strip("`")
            data = json.loads(raw)
            if isinstance(data, list):
                return [{
                    "question":       str(d.get("question", "")),
                    "option_a":       str(d.get("option_a", "")),
                    "option_b":       str(d.get("option_b", "")),
                    "option_c":       str(d.get("option_c", "")),
                    "option_d":       str(d.get("option_d", "")),
                    "correct_answer": str(d.get("correct_answer", "A")).upper()[0],
                    "explanation":    str(d.get("explanation", "")),
                } for d in data if d.get("question")][:count]
        except Exception as e:
            print(f"MCQ error: {e}")
        return []

    # ── Short questions ────────────────────────────────────────
    def generate_short_questions(self, text, count=5):
        if not self.client:
            return []
        try:
            resp = self.client.chat.completions.create(
                messages=[
                    {"role": "system",
                     "content": "Create short-answer questions. Respond ONLY with a valid JSON array of objects with \"question\" and \"answer\" keys. No markdown."},
                    {"role": "user",
                     "content": f"Create exactly {count} questions from:\\n\\n{text[:2500]}"},
                ],
                model="llama-3.3-70b-versatile",
                temperature=0.7, max_tokens=900,
            )
            raw = re.sub(
                r"```[a-z]*", "", resp.choices[0].message.content).strip().strip("`")
            data = json.loads(raw)
            if isinstance(data, list):
                return [{"question": str(d.get("question", "")), "answer": str(d.get("answer", ""))}
                        for d in data if d.get("question") and d.get("answer")][:count]
        except Exception as e:
            print(f"Short question error: {e}")
        return []
